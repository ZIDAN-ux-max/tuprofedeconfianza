# -*- coding: utf-8 -*-
"""Funciones utilitarias generales: seguridad, lectura de PDF y niveles."""
import hashlib
import io
import PyPDF2
from pptx import Presentation
from datetime import datetime
from zoneinfo import ZoneInfo

ZONA_HORARIA = ZoneInfo("America/Lima")


def ahora_peru():
    """Fecha y hora actual en horario de Peru (UTC-5), sin importar en que
    zona horaria este el servidor donde corre la app (Streamlit Cloud usa
    UTC por defecto). Usar SIEMPRE esto en vez de datetime.now()."""
    return datetime.now(ZONA_HORARIA)


def hoy_peru():
    """Solo la fecha (sin hora) de hoy en horario de Peru. Usar SIEMPRE
    esto en vez de date.today()."""
    return ahora_peru().date()


def hash_password(password):
    """Genera un hash SHA-256 de la contrasena."""
    return hashlib.sha256(password.encode()).hexdigest()


def hash_texto(texto):
    """Genera un hash SHA-256 del contenido de texto de un documento, para
    detectar duplicados aunque el archivo se haya subido con otro nombre."""
    return hashlib.sha256(texto.strip().encode()).hexdigest()


def extraer_texto_pdf(archivo, max_caracteres=3000):
    """Extrae texto de un PDF. Por defecto solo 3000 caracteres (uso rapido
    en el chat), pero la biblioteca de documentos pide mucho mas
    (ver documentos.py) para no perder la mayor parte del contenido."""
    try:
        pdf_reader = PyPDF2.PdfReader(io.BytesIO(archivo.read()))
        texto = ""
        for page in pdf_reader.pages:
            texto += page.extract_text()
        return texto[:max_caracteres]
    except Exception:
        return None


def extraer_texto_pptx(archivo, max_caracteres=3000):
    """Extrae el texto de todas las diapositivas de un PPTX (titulos,
    contenido de cajas de texto y tablas), mismo uso que extraer_texto_pdf."""
    try:
        presentacion = Presentation(io.BytesIO(archivo.read()))
        texto = ""
        for diapositiva in presentacion.slides:
            for forma in diapositiva.shapes:
                if forma.has_text_frame:
                    for parrafo in forma.text_frame.paragraphs:
                        for run in parrafo.runs:
                            texto += run.text + " "
                    texto += "\n"
                if forma.has_table:
                    for fila in forma.table.rows:
                        for celda in fila.cells:
                            texto += celda.text + " "
                    texto += "\n"
        return texto[:max_caracteres]
    except Exception:
        return None


def obtener_nivel(total):
    """Devuelve el nivel (texto + color) segun el total de preguntas hechas."""
    if total < 10:
        return "Principiante 🌱", "#6B7280"
    elif total < 50:
        return "Intermedio ⭐", "#D97706"
    elif total < 100:
        return "Avanzado 🔥", "#2563EB"
    else:
        return "Experto 👑", "#7C3AED"


def dividir_en_fragmentos(texto, tamano=900, solape=150):
    """Divide un texto largo (ej: un PDF de varias paginas) en fragmentos
    pequenos con algo de solape entre ellos, para poder buscar despues solo
    los fragmentos relevantes a una pregunta en vez de mandar todo el texto
    junto al tutor."""
    fragmentos = []
    inicio = 0
    n = len(texto)
    while inicio < n:
        fin = min(inicio + tamano, n)
        fragmentos.append(texto[inicio:fin])
        if fin == n:
            break
        inicio = fin - solape
    return fragmentos


def normalizar_latex(texto):
    """La IA a veces escribe formulas con \\[ \\] o \\( \\) en vez de $$ $$
    o $ $ (que es lo que el renderizador de Streamlit reconoce). Tambien a
    veces usa corchetes simples [ ] sin barra invertida, en su propia linea,
    como forma de 'destacar' una formula. Esto convierte todo eso al formato
    que Streamlit si reconoce, sin depender 100% de que la IA siga el
    formato pedido al pie de la letra."""
    if not texto:
        return texto
    texto = texto.replace("\\[", "$$").replace("\\]", "$$")
    texto = texto.replace("\\(", "$").replace("\\)", "$")

    import re

    def _convertir_linea(match):
        contenido = match.group(1).strip()
        return f"$${contenido}$$"

    # lineas que son SOLO una formula entre corchetes simples (ej: "[ v(t)=3t+1 ]")
    texto = re.sub(r'(?m)^\s*\[\s*(.+?)\s*\]\s*$', _convertir_linea, texto)

    # bloques \begin{aligned}...\end{aligned} (o align/cases/matrix) que la IA
    # a veces manda sueltos, sin envolver en $$ $$. Si ya vienen envueltos en
    # $$, el (?<!\$) / (?!\$) evita tocarlos de nuevo.
    texto = re.sub(
        r'(?<!\$)\\begin\{(aligned|align|cases|matrix|pmatrix|bmatrix)\}.*?\\end\{\1\}(?!\$)',
        lambda m: f"$${m.group(0)}$$",
        texto,
        flags=re.DOTALL
    )

    # si una formula (con simbolo $) quedo atrapada dentro de un <h4> o <li>,
    # Streamlit no la renderiza ahi adentro. Quitamos esa etiqueta especifica
    # (solo cuando tiene una formula) para que el LaTeX si se muestre bien.
    texto = re.sub(r'<h4[^>]*>(.*?\$.*?)</h4>', lambda m: f"\n**{m.group(1)}**\n", texto, flags=re.DOTALL)
    texto = re.sub(r'<li>(.*?\$.*?)</li>', lambda m: f"- {m.group(1)}\n", texto, flags=re.DOTALL)

    return texto
